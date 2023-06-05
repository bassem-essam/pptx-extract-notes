import collections 
import collections.abc
from pptx import Presentation
import json
from glob import  glob


template  = """
    <!DOCTYPE html>
    <html lang="en">
    <head>
        <meta charset="UTF-8" />
        <meta http-equiv="X-UA-Compatible" content="IE=edge" />
        <meta name="viewport" content="width=device-width, initial-scale=1.0" />
        <style>
        html,
        body {
            padding: 0;
            margin: 0;
        }

        #nav {
            display: flex;
            flex-direction: row;
            align-items: center;
            justify-content: center;
            background-color: #000;
            color: #fff;
            margin: 0 0 20px 0;
        }

        #nav > * {
            margin: 10px;
        }

        #text {
            text-indent: 2em;
            padding: 25px;
        }

        #slide-number {
            width: 4em;
        }
        </style>
        <title>NotesView</title>
    </head>
    <body>
        <div id="nav">
        <button id="btn-home">Home</button>
        <button id="btn-prev">Previous</button>
        <h3>Slide <input type="number" id="slide-number"></input></h3>
        <button id="btn-go">Go</button>
        <button id="btn-next">Next</button>
        <button id="btn-end">End</button>
        </div>
        <div id="text">PLACEHOLDER</div>
        <script>var data = FILL_POINT</script>
        <script>
        let number = 0;
        let slideNumber = document.getElementById("slide-number")

        document
            .getElementById("btn-go")
            .addEventListener("click", function (e) {
            n = Number(slideNumber.value);
            number = constrain(n - 1);
            Show();
            });

            slideNumber.addEventListener("keydown", function (e) {
                if (e.key == "Enter") {
                    n = Number(slideNumber.value);
                    number = constrain(n - 1);
                    Show();
                }
            })

        document
            .getElementById("btn-prev")
            .addEventListener("click", function (e) {
            number = constrain(number - 1);
            Show();
            });

        document
            .getElementById("btn-next")
            .addEventListener("click", function (e) {
            number = constrain(number + 1);
            Show();
            });

        document
            .getElementById("btn-home")
            .addEventListener("click", function (e) {
            number = lowest;
            Show();
            });

        document
            .getElementById("btn-end")
            .addEventListener("click", function (e) {
            number = highest;
            Show();
            });

        document.addEventListener("keydown", function (e) {
            if (e.key == "ArrowRight" || e.key == "PageDown") {
            number = constrain(number + 1);
            Show();
            } else if (e.key == "ArrowLeft" || e.key == "PageUp") {
            number = constrain(number - 1);
            Show();
            }
        });

        var lowest = 0;
        var highest = data.length - 1;

        function constrain(num) {
            return Math.max(Math.min(num, highest), lowest);
        }

        function Show() {        
            document.getElementById("text").innerText = data[number];
            document.getElementById("slide-number").value = number + 1;
        }

        Show();
        </script>
    </body>
    </html>
    """


def handleFile(file):
    ppt=Presentation(file)

    notes = []

    for page, slide in enumerate(ppt.slides):
        # this is the notes that doesn't appear on the ppt slide,
        # but really the 'presenter' note. 
        textNote = slide.notes_slide.notes_text_frame.text
        notes.append(textNote) 

    ready = template.replace("FILL_POINT", json.dumps(notes))
    templateFile = open(file + ".html", "w")
    templateFile.write(ready)
    templateFile.close()


for f in glob("*.pptx"):
    print(f)
    handleFile(f)
